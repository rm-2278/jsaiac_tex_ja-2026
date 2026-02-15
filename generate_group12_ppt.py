from pathlib import Path

import pptx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent

ASSETS = {
    "subgoal_early": ROOT / "presentation-videos/subgoal-visualization/early-stage-96k.gif",
    "subgoal_mid": ROOT / "presentation-videos/subgoal-visualization/mid-stage-296k.gif",
    "subgoal_late": ROOT / "presentation-videos/subgoal-visualization/late-stage-395k.gif",
    "policy_early": ROOT / "presentation-videos/policy-behavior/early-policy-42k.gif",
    "policy_mid": ROOT / "presentation-videos/policy-behavior/mid-policy-170k.gif",
    "policy_late": ROOT / "presentation-videos/policy-behavior/late-policy-368k.gif",
    "timeline_early": ROOT / "presentation-videos/timeline-progression/very-early-10k.gif",
    "timeline_comp": ROOT / "presentation-videos/timeline-progression/comparison-232k.gif",
    "freeway_score": ROOT / "media/atari/atari_freeway-scores.png",
}


def add_bullet(text_frame, text: str, level: int = 0, font_size: int = 24, bold: bool = False):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)
    p.font.bold = bold
    return p


def add_placeholder_box(slide, left, top, width, height, label: str):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
    shape.line.color.rgb = RGBColor(100, 100, 100)
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.font.size = Pt(18)
    return shape


def add_picture_or_placeholder(slide, path: Path, left, top, width, height, fallback_label: str):
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    else:
        add_placeholder_box(slide, left, top, width, height, f"Missing asset:\n{fallback_label}")


def build_presentation(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "階層的世界モデル Hieros の\n実態評価と限界の可視化"
    slide.placeholders[1].text = "Group 12\n提出ファイル: group12_research_overview.mp4"

    # Slide 2: Background & Objective
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "背景と目的"

    tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.2), Inches(5.2))
    tf = tx_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "背景"
    p.font.bold = True
    p.font.size = Pt(30)
    add_bullet(tf, "世界モデルの利点：サンプル効率と長期計画", font_size=24)
    add_bullet(tf, "既存モデル：Dreamer / Director / Hieros", font_size=24)

    p = tf.add_paragraph()
    p.text = "\n目的"
    p.font.bold = True
    p.font.size = Pt(30)
    add_bullet(tf, "性能向上だけでなく『階層性が機能しているか』を検証", font_size=24)
    add_bullet(tf, "サブゴールや行動の中身を可視化して評価", font_size=24)

    add_placeholder_box(
        slide,
        Inches(7.2),
        Inches(2.0),
        Inches(5.3),
        Inches(4.0),
        "推奨素材\n背景整理図\n(Dreamer / Director / Hieros)",
    )

    # Slide 3: Pinpad Results
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Visual Pinpad：サブゴール可視化と評価"

    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.9))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ハイパーパラメータ依存が強く、長期シーケンスの一貫性に課題。"
    p.font.size = Pt(24)

    img_y = Inches(2.1)
    img_w = Inches(3.9)
    img_h = Inches(4.3)

    add_picture_or_placeholder(slide, ASSETS["subgoal_early"], Inches(0.4), img_y, img_w, img_h, "early-stage-96k.gif")
    add_picture_or_placeholder(slide, ASSETS["subgoal_mid"], Inches(4.7), img_y, img_w, img_h, "mid-stage-296k.gif")
    add_picture_or_placeholder(slide, ASSETS["subgoal_late"], Inches(9.0), img_y, img_w, img_h, "late-stage-395k.gif")

    for x, label in [
        (Inches(0.4), "Early Stage (96k)"),
        (Inches(4.7), "Mid Stage (296k)"),
        (Inches(9.0), "Late Stage (395k)"),
    ]:
        cap = slide.shapes.add_textbox(x, Inches(6.55), img_w, Inches(0.4))
        cp = cap.text_frame.paragraphs[0]
        cp.text = label
        cp.alignment = PP_ALIGN.CENTER
        cp.font.size = Pt(14)

    # Slide 4: Atari Results
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Atari：スコアと方策のギャップ"

    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.9))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = "高スコアでも単純動作の繰り返しが多く、階層的な意図が見出しにくい。"
    p.font.size = Pt(24)

    add_picture_or_placeholder(
        slide,
        ASSETS["freeway_score"],
        Inches(0.8),
        Inches(2.1),
        Inches(5.8),
        Inches(3.9),
        "atari_freeway-scores.png",
    )

    behavior_paths = [ASSETS["policy_early"], ASSETS["policy_mid"], ASSETS["policy_late"]]
    x_positions = [Inches(7.0), Inches(8.9), Inches(10.8)]
    for i, (bp, x) in enumerate(zip(behavior_paths, x_positions), start=1):
        add_picture_or_placeholder(slide, bp, x, Inches(2.3), Inches(1.7), Inches(3.5), f"policy {i}")

    cap = slide.shapes.add_textbox(Inches(7.0), Inches(6.05), Inches(5.5), Inches(0.4))
    cp = cap.text_frame.paragraphs[0]
    cp.text = "Policy behavior: 42k / 170k / 368k"
    cp.alignment = PP_ALIGN.CENTER
    cp.font.size = Pt(14)

    # Slide 5: Time Progression
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "時間発展による内部表現の変化"

    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.9))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = "性能向上と内部の階層的表現の成熟は必ずしも一致しない。"
    p.font.size = Pt(24)

    add_picture_or_placeholder(
        slide,
        ASSETS["timeline_early"],
        Inches(1.7),
        Inches(2.1),
        Inches(4.2),
        Inches(3.8),
        "very-early-10k.gif",
    )
    add_picture_or_placeholder(
        slide,
        ASSETS["timeline_comp"],
        Inches(7.2),
        Inches(2.1),
        Inches(4.2),
        Inches(3.8),
        "comparison-232k.gif",
    )

    # Slide 6: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "結論と今後の課題"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "1. Visual Pinpadではハイパーパラメータ依存が強く、頑健性が低い"
    p.font.size = Pt(26)

    add_bullet(tf, "2. Atariでは高スコアでも単純方策が多く、期待した階層性が確認しづらい", font_size=26)
    add_bullet(tf, "3. 階層数増加で学習安定性が低下する", font_size=26)
    add_bullet(tf, "", font_size=12)
    last = add_bullet(tf, "今後の展望：動的抽象化メカニズムの改善と安定学習手法の確立", font_size=26, bold=True)
    last.font.color.rgb = RGBColor(0, 112, 192)

    prs.save(str(output_path))


if __name__ == "__main__":
    out = ROOT / "group12_research_overview.pptx"
    build_presentation(out)
    print(f"Created: {out}")
